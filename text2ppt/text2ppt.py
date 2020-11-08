#!/usr/bin/env python
# -*- coding: utf-8 -*-

import json
import csv
import os
import shutil

from pptx import Presentation

DIR = 'C:/dev/text2ucgppt/resources'
FILENAME = 'UCG-Hymnal-Lyrics-windows.pptx'
FILENAME_TEST = 'UCG-Hymnal-Lyrics-windows_TEST.pptx'

ABS_PATH = DIR+'/'+FILENAME

HYMN_SLIDE_JSON_PATH = '../resources/hymn_slide.json'

TITLE = 0
VERSE_SLIDE = 1
VERSE_SLIDE_END = 2
BLANK_SLIDE = 3
EMPTY_SLIDE = 4

NUM_OF_SLIDES = 2878

HYMN_TEXT_PATH = '../resources/hymn_text.csv'


def load_hymn_text():
    with open(HYMN_TEXT_PATH, mode='r', newline='', encoding='utf-8') as csv_file:
        csv_reader = csv.reader(csv_file, delimiter=':')
        data = list(csv_reader)

    #for row in data[1:]:
     #   print(row)

    return data[1:]


with open(HYMN_SLIDE_JSON_PATH) as json_file:
    hymn_slide_json = json.load(json_file)

print(hymn_slide_json)

prs = Presentation(ABS_PATH)

hymn_text = load_hymn_text()
hymn_number = int(hymn_text[0][0])

# save hymn_text file
DEST_DIR = '../resources/saved_hymn_texts'
shutil.copy(HYMN_TEXT_PATH, DEST_DIR)
dst_file = os.path.join(DEST_DIR, 'hymn_text.csv')
new_name = 'hymn_text_'+str(hymn_number)+'.csv'
new_dst_file_name = os.path.join(DEST_DIR, new_name)
os.rename(dst_file, new_dst_file_name)  #rename

# get start and end slide numbers
start = hymn_slide_json[hymn_number - 1]['slide']  # [0][0] because it's a list of a list
print(start)
end = hymn_slide_json[hymn_number]['slide']
print(end)

slides = prs.slides

# replace text
i = 0
for slide_index in range(start-1, end-2):  # end-2 because the slide between hymns is empty and would result in out of range error
    slide = slides[slide_index]
    print(slide_index)
    i = i+1
    #slide = slides[0]
    shape = slide.shapes[0] # always the 0th shape
    if shape.has_text_frame:
        text_frame = shape.text_frame
        print(len(text_frame.paragraphs))
        text = text_frame.paragraphs[0].text # only ever one paragraph
        print('FANCY TEXT: ', text)

        shape.text = hymn_text[i][0]  # more silliness with list in lists
        text = text_frame.paragraphs[0].text
        print('FANCY TEXT neu: ', text)
    else:
        print('NO TEXT')


if len(slides) != NUM_OF_SLIDES:
    print('!!! WARNING !!! Number of slides changed, not saving!!! Current: {}, Should: {}'.format(len(slides), NUM_OF_SLIDES))
else:
    prs.save(ABS_PATH)



